import tkinter as tk
from tkinter import *
from tkinter import filedialog
import os
from PIL import Image,ImageTk
from functools import partial
from tkinter import messagebox
import qa_main
import topic_find
import summ
import json
import random
import nltk
from fuzzywuzzy import process
from nltk.corpus import stopwords 
from nltk.tokenize import sent_tokenize, word_tokenize
from pptx import Presentation
import os, glob

#variables
file_clicked=False
search_clicked=False
material=""
word="preprocessor"
book1=""
book2=""
book3=""
book4=""
book5=""
book1_add=""
book2_add=""
book3_add=""
book4_add=""
book5_add=""
user_answer=[]
global f3there
f3there=1
global quizgenerated
quizgenerated=False


#functions

def searchslide(event):
    global search_clicked
    file_close()
    if(search_clicked==True):
        search_clicked=False
        print("it should be gray")
        lbl11.bind("<Leave>", partial(imgcolor_config, lbl11, search_gray))
        f4.grid_remove()
    else:
        search_clicked=True
        lbl11.bind("<Leave>", partial(imgcolor_config, lbl11, search_white))
        print("it should be white")
        f4.grid()

def search_close():
    global search_clicked
    search_clicked=False
    print("it should be gray")
    #lbl11.bind("<Leave>", partial(imgcolor_config, lbl11, search_gray))
    lbl11.configure(image=search_gray)
    f4.grid_remove()

def fileslide(event):
    global file_clicked
    search_close()
    if(file_clicked==True):
        file_clicked=False
        print("it should be gray")
        lbl10.bind("<Leave>", partial(imgcolor_config, lbl10, file_gray))
        f4.grid_remove()
    else:
        file_clicked=True
        lbl10.bind("<Leave>", partial(imgcolor_config, lbl10, file_white))
        print("it should be white")
        f4.grid()

def file_close():
    global file_clicked
    file_clicked=False
    print("it should be gray")
    #lbl10.bind("<Leave>", partial(imgcolor_config, lbl10, file_gray))
    lbl10.configure(image=file_gray)
    f4.grid_remove()

def worked(event):
    read_file(book1_add)

def work(event):

    print("it worked")
    gradient.filename=filedialog.askopenfilename(initialdir="C:/Users/manda/Desktop/texts",title="select file",filetypes=(("txt","*.txt"),("pdf","*.pdf"),("document","*.doc")))
    if(gradient.filename):
        try:
            file_read(gradient.filename)
            global book1_add
            book1_add = (gradient.filename)
        except:
            messagebox.showerror("Error", "can't open file")

def file_read(fname):
    txt = open(fname)
    material=txt.read()
    print(material)
    lbl9.configure(font="times 12 bold")
    lbl9.tag_configure(gradient,justify='center')
    lbl9.delete('1.0',END)
    lbl9.insert('1.0',material)
    topic_find.filen=fname
    os.system("python topic_find.py")
    book1=topic_find.top_3_title[1]
    lbl1.configure(text=book1,font="times 15 bold")

def read_file(fname):
    txt = open(fname)
    material=txt.read()
    print(material)
    lbl9.configure(font="times 12 bold")
    #lbl9.clear()
    lbl9.delete('1.0',END)
    lbl9.insert('1.0',material)

def addqa():
    
    
    x=qasea.get()
    print(x)
    lbl9.insert(END,"\n Q. "+x)
    qasea.delete(0,END)
    #import numpy
    stop_words = set(stopwords.words('english')) 
    h=0
    #index=-1
    #mindex=0
    match="Sorry i dont understand your question"
    w="0"
    txt = open(book1_add)
    material=txt.read()
    s=material
    #s="Compiler is as software which converts a program written in high level language to low level language.Compiler converts the program into machine level language. There is another compiler known as Cross-Compiler. Cross-Compiler runs on one machine but generates machine or assembly code for another machine. There are six phases in a compiler, they are Lexical Analysis, Syntax Analysis, Semantic Analysis, Intermediate Code Generation, Code Optimization, Code Generation. Phases of Compilation are divided into two groups, Front-end and Back-end. Front-end consists of Lexical Analysis, Syntax Analysis, Semantic Analysis and Intermediate Code Generation. Back-end consists of Code Optimization, Code Generation and Assembly phase. These phases ensure that the speed of execution is faster in the compiler."
    st=sent_tokenize(s)
    #q=input("Enter your question: ")
    q=x
    stf=""
    l=0
    for ans in st:
        i = nltk.word_tokenize(ans)
        #si= [x for x in i if not x in stop_words]
        t = nltk.pos_tag(i) 
        #ne= nltk.ne_chunk(t) 
        for gram in t:
            if gram[1]=="NNP":
                w=gram[0]
                stf+=" "+w
            elif gram[1]==".":
                stf+=gram[0]
            elif gram[1]=="PRP" and w!="0":
                stf+=" "+w
            else:
                if l==0:
                    stf+=gram[0]
                else:
                    stf+=" "+gram[0]
                l=1
    stf=sent_tokenize(stf)
    Ratios = process.extract(q,st)
    print(Ratios)
    for w in Ratios:
        #index=index+1
        if w[1] >=h and w[1]>30:
            h=w[1]
            match = w[0]
            #mindex=index
    #print(match)
    lbl9.insert(END,"\n Ans. "+match)
    lbl9.see(END)

def questionanswer(event):
    qa_main.inputTextPath=gradient.filename
    os.system("python qa_main.py")
    file_read("DB/qa_output.txt")
    f3.grid()
    global f3there
    f3there=0
    if(quizgenerated==True):
        quizbut.grid_remove()
        quizlbl.grid_remove()
        quizradio.grid_remove()


def incrementi_and_much_more(g):

  print("here")
  with open('DB/quiz_output.txt') as f:
        data = json.load(f)
        
  if(data[g]['similar_words']):
      print(user_answer)
  else:
      x=fib.get()
      print(x)
      user_answer.append(x)
      print(user_answer)
      fib.delete(0,END)

  g=g+1
  if(g>0):
    generate_quiz(g)

def show_result(g):
    print("rsult")
    with open('DB/quiz_output.txt') as f:
        data = json.load(f)
    if(data[g]['similar_words']):
      print(user_answer)
    else:
      x=fib.get()
      print(x)
      user_answer.append(x)
      print(user_answer)
      fib.delete(0,END)
    ran=len(data)
    score=0
    for i in range(0,ran):
        if(user_answer[i]==data[i]['answer']):
            score=score+1

    tk.messagebox.showinfo(title="result", message="you scored"+str(score)+"/"+str(ran))

def selected1():
    #r1.select()
    #r1.configure(fg="white",selectcolor="white")
    x = radiovar.get()
    print(x)
    user_answer.append(x)
    

def selected2():
    #r2.configure(selectcolor="yellow")
    x = radiovar.get()
    print(x)
    user_answer.append(x)
    

def selected3():
    #r3.select()
    x = radiovar.get()
    print(x)
    user_answer.append(x)
    

def selected4():
    #r4.select()
    x = radiovar.get()
    print(x)
    user_answer.append(x)
    

def generate_quiz(g):
    global quizgenerated
    quizgenerated=True

    if(f3there==0):
        f3.grid_remove()
    lbl9.delete('1.0',END)
    quizlbl.grid_propagate(False)
    quizlbl.grid(column=0,row=0,sticky='N'+'S'+'W'+'E')
    quizradio.grid(column=0,row=1,sticky='N'+'S'+'W'+'E')
    quizbut.grid(column=0,row=2,sticky='N'+'S'+'W'+'E')
    with open('DB/quiz_output.txt') as f:
        data = json.load(f)
    ran=len(data)
    print(ran)
    if(g<ran-1):
        whichbut=1
        qq.pack_propagate(False)

        qq.pack(pady=50)
        
        qq.configure(text=data[g]['question'])
        #lbl9.insert(1.0,data[i]['question']+"\n")
        if(data[g]['similar_words']):
            for j in range(1,2):
                #print(data[i]['similar_words'][j])
                rannn=random.randrange(1,5)
                global radiovar
    
                radiovar.set("1")
                
                #print(ran)
                if(rannn==1):
                    #print(data[i]['similar_words'][j])
                    r1.pack()
                    r1.configure(text=data[g]['answer'],value=data[g]['answer'])
                    r2.pack()
                    r2.configure(text=data[g]['similar_words'][0],value=data[g]['similar_words'][0])
                    r3.pack()
                    r3.configure(text=data[g]['similar_words'][1],value=data[g]['similar_words'][1])
                    r4.pack()
                    r4.configure(text=data[g]['similar_words'][2],value=data[g]['similar_words'][2])

                elif(rannn==2):
                    #print(data[i]['similar_words'][j])
                    r1.pack()
                    r1.configure(text=data[g]['similar_words'][0],value=data[g]['similar_words'][0])
                    r2.pack()
                    r2.configure(text=data[g]['answer'],value=data[g]['answer'])
                    r3.pack()
                    r3.configure(text=data[g]['similar_words'][1],value=data[g]['similar_words'][1])
                    r4.pack()
                    r4.configure(text=data[g]['similar_words'][2],value=data[g]['similar_words'][2])
                elif(rannn==3):
                    # print(data[i]['similar_words'][j])
                    r1.pack()
                    r1.configure(text=data[g]['similar_words'][0],value=data[g]['similar_words'][0])
                    r2.pack()
                    r3.configure(text=data[g]['similar_words'][1],value=data[g]['similar_words'][1])
                    r3.pack()
                    r3.configure(text=data[g]['answer'],value=data[g]['answer'])
                    r4.pack()
                    r4.configure(text=data[g]['similar_words'][2],value=data[g]['similar_words'][2])
                else:
                    #print(data[i]['similar_words'][j])
                    r1.pack()
                    r1.configure(text=data[g]['similar_words'][0],value=data[g]['similar_words'][0])
                    r2.pack()
                    r2.configure(text=data[g]['similar_words'][1],value=data[g]['similar_words'][1])
                    r3.pack()
                    r3.configure(text=data[g]['similar_words'][2],value=data[g]['similar_words'][2])
                    r4.pack()
                    r4.configure(text=data[g]['answer'],value=data[g]['answer'])
        else:
            #print("fill in the blanck")
            r1.destroy()
            r2.destroy()
            r3.destroy()
            r4.destroy()
            fib.pack(side=TOP)

        nextbut.configure(command=lambda :incrementi_and_much_more(g))
        nextbut.pack()
        print(g)

    else:
        whichbut=0
    
        
    if(whichbut==1):
        nextbut.configure(command=lambda :incrementi_and_much_more(g))
        nextbut.pack()
        print(g)
    else:
        submitbut.configure(command=lambda :show_result(g))
        nextbut.destroy()
        submitbut.pack()
        print(g)
        

        
        


def quiz(event):
    os.system("wikitrivia --output DB/quiz_output.txt "+word)
    #file_read("DB/quiz_output.txt")
    generate_quiz(0)

def summm(event):
    
    summ.filen=gradient.filename
    os.system("python summ.py")
    file_read("DB/summ_output.txt")
    
    #this part generates a ppt but is not comlpete
    '''
    if(f3there==0):
        f3.grid_remove()
    if(quizgenerated==True):
        quizbut.grid_remove()
        quizlbl.grid_remove()
        quizradio.grid_remove() 



    prs = Presentation()

    def slide(number=0):
        title_slide_layout = prs.slide_layouts[number]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = file.readline()
        file.readline()
        subtitle.text = file.readline()
        file.readline() 


    file_chosen = "DB/hehe.txt"
    with open(file_chosen, 'r', encoding="utf-8") as file:
        file_lenght = len(file.readlines())
        file.seek(0)
        slide(0)
        for i in range(int(file_lenght / 4)):
            slide(1)
    file_to_save = file_chosen[:-4] + ".pptx"
    prs.save("hehe.pptx")
    os.startfile("hehe.pptx")   
    '''

def color_config(widget, color, event):
    widget.configure(foreground=color)

def imgcolor_config(widget, color, event):
    widget.configure(image=color)



root=tk.Tk()
root.geometry("1100x600")
#root.resizable(0,0)
screenw=root.winfo_screenwidth()
screenh=root.winfo_screenheight()

#title
root.title("Bibliophile Companion")

#icon
photo = PhotoImage(file = "icons/book.gif")
root.iconphoto(False, photo)

#menu
w = Menu(root)  
file=Menu(w,tearoff=0)
file.add_command(label="open")
w.add_cascade(label="file",menu=file)
w.add_command(label="edit")
w.add_command(label="faq")

#partioning
p = PanedWindow(root,width=1100,height=600, orient=HORIZONTAL,bg="black",bd=2)
p.pack(padx=4,pady=3,expand=1,fill=BOTH)
f1 = LabelFrame(p, text='',height=600,width=70,bg="#181818",fg="white",bd=0,labelanchor="n")
f2 = LabelFrame(p, text='',height=600,width=140,bg="#1D191B",fg="white",bd=0)
f4 = LabelFrame(p, text='',height=600,width=950,bg="#2b2b2b",fg="white",bd=0,labelanchor="n")  

Grid.rowconfigure(p, 0, weight=1)
Grid.columnconfigure(p, 2, weight=1)

f1.grid(column=0,row=0,sticky='N'+'S'+'W')
f4.grid(column=1,row=0,sticky='N'+'S'+'W'+'E')
f2.grid(column=2,row=0,sticky='N'+'S'+'W'+'E')
f4.grid_remove()

gradient=LabelFrame(f2,height=500,width=950,bg="black")
f3=LabelFrame(f2,text="",height=100,width=950,bg="#2b2b2b",fg="white",bd=0,labelanchor="n")
Grid.rowconfigure(f2, 0, weight=1)
Grid.rowconfigure(f2, 1, weight=0)
Grid.columnconfigure(f2, 0, weight=1)
gradient.grid(row=0,column=0,sticky='N'+'S'+'W'+'E')
f3.grid(row=1,column=0,sticky='N'+'S'+'W'+'E')
f3.grid_remove()
gradient.grid_propagate(False)

qasea=Entry(f3)
qaseabut=Button(f3,text="search",command=lambda :addqa())
Grid.rowconfigure(f3, 0, weight=3)
Grid.rowconfigure(f3, 1, weight=1)
Grid.columnconfigure(f3, 0, weight=1)
qasea.grid(row=0,column=0,sticky='w'+'e')
qaseabut.grid(row=0,column=1)

functi=LabelFrame(f1, text='',height=500,width=70,bg="#181818",fg="white",bd=0,labelanchor="n")
addd=LabelFrame(f1, text='',height=100,width=70,bg="#181818",fg="white",bd=0,labelanchor="n")
Grid.rowconfigure(f1, 0, weight=6)
Grid.rowconfigure(f1, 1, weight=1)
Grid.columnconfigure(f1, 0, weight=1)
functi.grid(row=0,column=0,sticky='N'+'S'+'W'+'E')
addd.grid(row=1,column=0,sticky='N'+'S'+'W'+'E')
addd.grid_propagate(False)

#content

#f1
qa_white=ImageTk.PhotoImage(Image.open("icons/qa_white.png").resize((35,30)))
qa_gray=ImageTk.PhotoImage(Image.open("icons/qa_gray.png").resize((35,30)))
quiz_gray=ImageTk.PhotoImage(Image.open("icons/quiz_gray.png").resize((33,33)))
quiz_white=ImageTk.PhotoImage(Image.open("icons/quiz_white.png").resize((33,33)))
ppt_white=ImageTk.PhotoImage(Image.open("icons/ppt_white.png").resize((30,33)))
ppt_gray=ImageTk.PhotoImage(Image.open("icons/ppt_gray.png").resize((30,33)))
file_white=ImageTk.PhotoImage(Image.open("icons/file_white.png").resize((30,30)))
file_gray=ImageTk.PhotoImage(Image.open("icons/file_gray1.png").resize((30,30)))
search_white=ImageTk.PhotoImage(Image.open("icons/search_white.png").resize((30,30)))
search_gray=ImageTk.PhotoImage(Image.open("icons/search_gray.png").resize((30,30)))
lbl5=Label(functi,image=qa_gray,takefocus=True,bg="#181818",fg="gray",font="times 24 bold")
lbl5.bind("<Enter>", partial(imgcolor_config, lbl5, qa_white))
lbl5.bind("<Leave>", partial(imgcolor_config, lbl5, qa_gray))
lbl5.grid(row=2,padx=20,pady=20)
lbl5.bind("<Button-1>",questionanswer)
lbl6=Label(functi,image=quiz_gray,takefocus=True,bg="#181818",fg="gray",font="times 24 bold")
lbl6.bind("<Enter>", partial(imgcolor_config, lbl6, quiz_white))
lbl6.bind("<Leave>", partial(imgcolor_config, lbl6, quiz_gray))
lbl6.grid(row=3,padx=10,pady=20)
lbl6.bind("<Button-1>",quiz)
lbl7=Label(functi,image=ppt_gray,takefocus=True,bg="#181818",fg="gray",font="times 24 bold")
lbl7.bind("<Enter>", partial(imgcolor_config, lbl7, ppt_white))
lbl7.bind("<Leave>", partial(imgcolor_config, lbl7, ppt_gray))
lbl7.grid(row=4,padx=10,pady=20)
lbl7.bind("<Button-1>",summm)
lbl10=Label(functi,image=file_gray,takefocus=True,bg="#181818",fg="gray",font="times 24 bold")
lbl10.bind("<Enter>", partial(imgcolor_config, lbl10, file_white))
lbl10.bind("<Leave>", partial(imgcolor_config, lbl10, file_gray))
lbl10.grid(row=0,padx=10,pady=20)
lbl10.bind("<Button-1>",fileslide)
lbl11=Label(functi,image=search_gray,takefocus=True,bg="#181818",fg="gray",font="times 24 bold")
lbl11.bind("<Enter>", partial(imgcolor_config, lbl11, search_white))
lbl11.bind("<Leave>", partial(imgcolor_config, lbl11, search_gray))
lbl11.grid(row=1,padx=10,pady=20)
lbl11.bind("<Button-1>",searchslide)


addbut_white=ImageTk.PhotoImage(Image.open("icons/addbut_white.png").resize((50,50)))
addbut_gray=ImageTk.PhotoImage(Image.open("icons/addbut_gray.png").resize((50,50)))
lbl8=Label(addd,image=addbut_gray,takefocus=True,bg="#181818",fg="gray",font="times 24 bold")
lbl8.bind("<Enter>", partial(imgcolor_config, lbl8, addbut_white))
lbl8.bind("<Leave>", partial(imgcolor_config, lbl8, addbut_gray))
#lbl8.grid(column=1,row=1,padx=10,pady=20)
lbl8.grid(padx=10,pady=10,sticky='s')
lbl8.bind("<Button-1>",work)


#f2
lbl9=Text(gradient,bg="black",fg="gray",height=500,width=950,font="times 24 bold",spacing1=1,wrap=WORD,bd=0)
Grid.rowconfigure(gradient, 0, weight=1)
Grid.columnconfigure(gradient, 0, weight=1)
lbl9.grid(column=0,row=0,sticky='N'+'S'+'W'+'E',padx=5,pady=5)
lbl9.grid_propagate(False)

#quizlbl
quizlbl=LabelFrame(lbl9,bg="black",fg="gray",height=100,width=950,font="times 24 bold",bd=0)
qq=Label(quizlbl,text="",bg="black",fg="gray",wraplength=600,font="times 12 bold")
quizradio=LabelFrame(lbl9,bg="black",fg="gray",height=300,width=950,font="times 24 bold",bd=0)
radiovar=StringVar(value="1")
#r1=Radiobutton(quizradio,text="12",variable = radiovar,value=1,bg="black",fg="white",selectcolor="black",activebackground="black",activeforeground="white",command = lambda :selected1())
r1=Radiobutton(quizradio,text="",bg="black",fg="white",pady=10,justify=LEFT,value=1,variable = radiovar,activebackground="black",activeforeground="white",command = lambda :selected1(),selectcolor="black")
r2=Radiobutton(quizradio,text="",bg="black",fg="white",pady=10,justify=LEFT,value=2,variable = radiovar,activebackground="black",activeforeground="white",command = lambda :selected2(),selectcolor="black")
r3=Radiobutton(quizradio,text="",bg="black",fg="white",pady=10,justify=LEFT,value=3,variable = radiovar,activebackground="black",activeforeground="white",command = lambda :selected3(),selectcolor="black")
r4=Radiobutton(quizradio,text="",bg="black",fg="white",pady=10,justify=LEFT,value=4,variable = radiovar,activebackground="black",activeforeground="white",command = lambda :selected4(),selectcolor="black")
quizbut=LabelFrame(lbl9,bg="black",fg="gray",height=100,width=950,font="times 24 bold",bd=0)
Grid.rowconfigure(lbl9, 0, weight=1)
Grid.rowconfigure(lbl9, 1, weight=3)
Grid.rowconfigure(lbl9, 2, weight=1)
Grid.columnconfigure(lbl9, 0, weight=1)
fib=Entry(quizradio)

nextbut=Button(quizbut,text='next',height=2,width=7,bg="white")
submitbut=Button(quizbut,text='submit',height=2,width=7,bg="white")

#f4
lbl1=Label(f4,text=book1,takefocus=True,pady=20,bg="#2b2b2b",fg="gray",font="times 24 bold")
lbl1.bind("<Enter>", partial(color_config, lbl1, "#ffffff"))
lbl1.bind("<Leave>", partial(color_config, lbl1, "gray"))
lbl1.grid(row=0)
lbl1.bind("<Button-1>",worked)
lbl2=Label(f4,text=book2,takefocus=True,pady=20,bg="#2b2b2b",fg="gray",font="times 24 bold")
lbl2.bind("<Enter>", partial(color_config, lbl2, "#ffffff"))
lbl2.bind("<Leave>", partial(color_config, lbl2, "gray"))
lbl2.grid(row=1)
lbl2.bind("<Button-1>",worked)
lbl3=Label(f4,text=book3,takefocus=True,pady=20,bg="#2b2b2b",fg="gray",font="times 24 bold")
lbl3.bind("<Enter>", partial(color_config, lbl3, "#ffffff"))
lbl3.bind("<Leave>", partial(color_config, lbl3, "gray"))
lbl3.grid(row=2)
lbl3.bind("<Button-1>",worked)
lbl4=tk.Label(f4,text=book4,takefocus=True,pady=20,cursor='CIRCLE',padx=50,bg="#2b2b2b",fg="gray",font="times 24 bold")
lbl4.bind("<Enter>", partial(color_config, lbl4, "#ffffff"))
lbl4.bind("<Leave>", partial(color_config, lbl4, "gray"))
lbl4.grid(row=3)
lbl4.bind("<Button-1>",worked)

root.configure(menu=w)
root.mainloop()
